import os
import sys
import argparse
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

try:
    import win32com.client
except ImportError:
    raise ImportError("pywin32 が必要です。pip install pywin32 を実行してください。")


ITEMS = {
    1: "体重",
    2: "血糖値",
    3: "血圧・心拍数",
}


def get_onedrive_path(*parts):
    return Path(os.environ["OneDrive"]).joinpath(*parts)


BASE_HEALTH = get_onedrive_path("ドキュメント", "PythonWork", "Health")
PNG_DIR = BASE_HEALTH / "health_graph_png"
PPTX_PATH = BASE_HEALTH / "健康管理グラフ.pptx"

EXCEL_PATH = get_onedrive_path("ドキュメント", "PythonWork", "ExcelDATA", "データ表1.xlsx")

GRAPH_SCRIPT = Path(__file__).parent / "04_charts" / "Graph_from_Excel(Complete_2025_11_03).py"


def run_graph_script(start_date, end_date, item):
    cmd = [
        sys.executable,
        str(GRAPH_SCRIPT),
        "--start-date", start_date,
        "--end-date", end_date,
        "--item", str(item),
    ]

    print("グラフ作成:", ITEMS[item])
    result = subprocess.run(cmd, text=True, capture_output=True)

    print(result.stdout)

    if result.returncode != 0:
        print(result.stderr)
        raise RuntimeError(f"グラフ作成に失敗しました: item={item}")


def export_excel_chart_to_png(item):
    PNG_DIR.mkdir(parents=True, exist_ok=True)

    png_path = PNG_DIR / f"{item}_{ITEMS[item]}.png"

    excel = win32com.client.DispatchEx("Excel.Application")
    excel.Visible = False
    excel.DisplayAlerts = False

    try:
        wb = excel.Workbooks.Open(str(EXCEL_PATH))
        ws = wb.Worksheets("Sheet3")

        charts = ws.ChartObjects()
        if charts.Count == 0:
            raise RuntimeError("Sheet3 にグラフがありません")

        chart = charts.Item(1).Chart
        chart.Export(str(png_path))

        wb.Close(SaveChanges=False)

    finally:
        excel.Quit()

    print("PNG保存:", png_path)
    return png_path


def add_slide(prs, title, png_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True

    slide.shapes.add_picture(
        str(png_path),
        Inches(0.5),
        Inches(0.9),
        width=Inches(9.0),
    )


def make_pptx(png_files, start_date, end_date):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for item, png_path in png_files:
        title = f"{ITEMS[item]}（{start_date} ～ {end_date}）"
        add_slide(prs, title, png_path)

    prs.save(PPTX_PATH)
    print("PowerPoint作成完了:", PPTX_PATH)

    return PPTX_PATH


def main():
    try:
        parser = argparse.ArgumentParser()
        parser.add_argument("--start-date", required=True)
        parser.add_argument("--end-date", required=True)
        args = parser.parse_args()

        png_files = []

        for item in [1, 2, 3]:
            run_graph_script(args.start_date, args.end_date, item)
            png_path = export_excel_chart_to_png(item)
            png_files.append((item, png_path))

        pptx_path = make_pptx(png_files, args.start_date, args.end_date)
        os.startfile(pptx_path)

    except Exception as e:
        import traceback

        log = BASE_HEALTH / "make_health_pptx_error.txt"
        with open(log, "w", encoding="utf-8") as f:
            traceback.print_exc(file=f)

        raise


if __name__ == "__main__":
    main()
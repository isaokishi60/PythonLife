from pathlib import Path
import os
import platform
import subprocess

from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

# ユーザー名 spax2 / kawgu の違いを自動吸収
BASE_DIR = Path.home() / "OneDrive" / "ドキュメント" / "PythonWork" / "Health"

PERIOD_PNG_DIR = BASE_DIR / "01_Garmin_Import" / "outputs_period" / "png"
NIGHT_PNG_DIR = BASE_DIR / "01_Garmin_Import" / "outputs" / "png"

OUT_PPTX = BASE_DIR / "心房細動グラフ.pptx"

# 16:9 ワイド画面
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 画像を入れる枠。Night Heart Rate もこの枠内に必ず収める。
IMG_BOX_X = Inches(0.35)
IMG_BOX_Y = Inches(0.55)
IMG_BOX_W = Inches(12.65)
IMG_BOX_H = Inches(6.70)

# スライド上のタイトル位置。
# 1～3枚目は、グラフ内の注記と重なりにくいよう左上・凡例の上に置く。
TITLE_X = Inches(0.55)
TITLE_Y = Inches(0.62)
TITLE_W = Inches(3.4)
TITLE_H = Inches(0.35)


def latest_file(folder: Path, pattern: str) -> Path:
    files = list(folder.glob(pattern))
    if not files:
        raise FileNotFoundError(f"ファイルが見つかりません: {folder}\\{pattern}")
    return max(files, key=lambda p: p.stat().st_mtime)


def add_picture_fit(slide, image_path: Path, box_x, box_y, box_w, box_h):
    """画像の縦横比を維持したまま、指定した枠内に収めて中央配置する。"""
    with Image.open(image_path) as img:
        px_w, px_h = img.size

    img_ratio = px_w / px_h
    box_ratio = box_w / box_h

    if img_ratio >= box_ratio:
        # 横長画像：幅を枠に合わせ、高さを縮小
        draw_w = box_w
        draw_h = int(box_w / img_ratio)
    else:
        # 縦長画像：高さを枠に合わせ、幅を縮小
        draw_h = box_h
        draw_w = int(box_h * img_ratio)

    draw_x = box_x + int((box_w - draw_w) / 2)
    draw_y = box_y + int((box_h - draw_h) / 2)

    return slide.shapes.add_picture(str(image_path), draw_x, draw_y, width=draw_w, height=draw_h)


def add_slide_title(slide, title: str):
    text_box = slide.shapes.add_textbox(TITLE_X, TITLE_Y, TITLE_W, TITLE_H)
    tf = text_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(16)
    run.font.bold = True
    return text_box


def open_pptx(path: Path):
    """作成した PowerPoint ファイルを最後に開く。Windows では既定アプリで開く。"""
    if platform.system() == "Windows":
        os.startfile(str(path))  # type: ignore[attr-defined]
    elif platform.system() == "Darwin":
        subprocess.Popen(["open", str(path)])
    else:
        subprocess.Popen(["xdg-open", str(path)])


slides = [
    ("DailyBeats", latest_file(PERIOD_PNG_DIR, "HeartPeriod_DailyBeats_2025-10-01_*.png")),
    ("RHR",        latest_file(PERIOD_PNG_DIR, "HeartPeriod_RHR_2025-10-01_*.png")),
    ("Tachy",      latest_file(PERIOD_PNG_DIR, "HeartPeriod_Tachy_2025-10-01_*.png")),
    ("NightHR",    latest_file(NIGHT_PNG_DIR, "RestHR_Night_*_21-06.png")),
]

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

for title, png in slides:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 画像は width だけ指定せず、縦横比を保って枠内に収める。
    add_picture_fit(slide, png, IMG_BOX_X, IMG_BOX_Y, IMG_BOX_W, IMG_BOX_H)

    # タイトルは画像の後に追加し、前面に出す。
    add_slide_title(slide, title)

    print(f"{title}: {png.name}")

prs.save(OUT_PPTX)
print(f"Saved: {OUT_PPTX}")
print(f"スライド数: {len(slides)}")

# 保存後、PowerPoint ファイルを開いた状態にする。
open_pptx(OUT_PPTX)
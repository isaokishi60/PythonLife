# %%
# HA1cと血糖値の平均日数分析
#
# 1. HA1c測定日前60日間の平均血糖値を計算
# 2. 60日平均の時系列グラフと回帰分析
# 3. 30日～120日を5日刻みで比較
# 4. 各日数の相関係数r、決定係数R²、回帰係数を計算
# 5. 最適日数を自動判定
# 6. Excelに分析結果を保存
# 7. 5枚のPowerPointを作成して自動表示

import os
from datetime import timedelta

import numpy as np
import pandas as pd
import openpyxl

import matplotlib.pyplot as plt
import matplotlib.dates as mdates

from openpyxl.chart import LineChart, Reference
from openpyxl.utils.dataframe import dataframe_to_rows

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


# ============================================================
# 基本設定
# ============================================================

STANDARD_DAYS = 60
AVERAGE_DAYS_LIST = list(range(30, 121, 5))

# 各平均期間について必要とする血糖値データの割合
# 例：60日平均では60×0.8=48件以上を必要とする
MIN_COVERAGE_RATIO = 0.80

ANALYSIS_SHEET_NAME = "平均日数分析"


# ============================================================
# ファイルパス
# ============================================================

def get_excel_path(filename, folder="ExcelDATA"):
    base = os.path.join(
        os.environ["OneDrive"],
        "ドキュメント",
        "PythonWork"
    )
    return os.path.join(base, folder, filename)


hba1c_path = get_excel_path(
    "データ表4.xlsx",
    folder="ExcelDATA"
)

glucose_path = get_excel_path(
    "データ表1.xlsx",
    folder="ExcelDATA"
)

output_dir = os.path.join(
    os.environ["OneDrive"],
    "ドキュメント",
    "PythonWork",
    "Health",
    "05_ha1c"
)

os.makedirs(output_dir, exist_ok=True)

time_series_png = os.path.join(
    output_dir,
    "HA1c_60日平均血糖値_時系列.png"
)

regression_60_png = os.path.join(
    output_dir,
    "HA1c_60日平均血糖値_回帰分析.png"
)

correlation_comparison_png = os.path.join(
    output_dir,
    "HA1c_平均日数別_相関係数.png"
)

r2_comparison_png = os.path.join(
    output_dir,
    "HA1c_平均日数別_決定係数.png"
)

prediction_png = os.path.join(
    output_dir,
    "血糖値から予測されるHA1c.png"
)

pptx_path = os.path.join(
    output_dir,
    "HA1cと平均血糖値_平均日数分析.pptx"
)


# ============================================================
# Matplotlibの日本語フォント
# ============================================================

plt.rcParams["font.family"] = [
    "Yu Gothic",
    "Meiryo",
    "MS Gothic",
    "sans-serif"
]
plt.rcParams["axes.unicode_minus"] = False


# ============================================================
# 共通関数
# ============================================================

def load_hba1c_data(excel_path):
    """
    データ表4.xlsxのSheet2から日付とHA1cを読み込む。

    初回の元シートがB列・C列の場合と、
    本プログラム実行後のA列「日付」・B列「HA1c」の場合の
    どちらにも対応する。
    """
    raw_df = pd.read_excel(
        excel_path,
        sheet_name="Sheet2"
    )

    normalized_columns = {
        str(col).strip(): col
        for col in raw_df.columns
    }

    if "日付" in normalized_columns and "HA1c" in normalized_columns:
        df = raw_df[
            [
                normalized_columns["日付"],
                normalized_columns["HA1c"]
            ]
        ].copy()
        df.columns = ["日付", "HA1c"]
    else:
        df = pd.read_excel(
            excel_path,
            sheet_name="Sheet2",
            usecols=[1, 2]
        )
        df.columns = ["日付", "HA1c"]

    df["日付"] = pd.to_datetime(
        df["日付"],
        errors="coerce"
    )

    df["HA1c"] = pd.to_numeric(
        df["HA1c"],
        errors="coerce"
    )

    df = (
        df.dropna(subset=["日付", "HA1c"])
        .sort_values("日付")
        .drop_duplicates(subset=["日付"], keep="last")
        .reset_index(drop=True)
    )

    return df


def load_glucose_data(excel_path):
    """
    データ表1.xlsxのSheet4から日付と血糖値を読み込む。
    """
    df = pd.read_excel(
        excel_path,
        sheet_name="Sheet4",
        usecols=[0, 2]
    )

    df.columns = ["日付", "血糖値"]

    df["日付"] = pd.to_datetime(
        df["日付"],
        errors="coerce"
    )

    df["血糖値"] = pd.to_numeric(
        df["血糖値"],
        errors="coerce"
    )

    df = (
        df.dropna(subset=["日付", "血糖値"])
        .sort_values("日付")
        .reset_index(drop=True)
    )

    return df


def calculate_average_glucose(
    hba1c_df,
    glucose_df,
    average_days,
    minimum_coverage_ratio=MIN_COVERAGE_RATIO
):
    """
    各HA1c測定日について、測定日前average_days日間の
    平均血糖値とデータ数を計算する。

    測定日当日は含めず、
    [測定日-average_days, 測定日) を対象とする。

    必要データ数に満たない場合、平均値はNaNとする。
    """
    average_list = []
    count_list = []

    minimum_count = int(
        np.ceil(average_days * minimum_coverage_ratio)
    )

    for hba1c_date in hba1c_df["日付"]:

        start_date = hba1c_date - timedelta(days=average_days)

        mask = (
            (glucose_df["日付"] >= start_date)
            & (glucose_df["日付"] < hba1c_date)
        )

        target = glucose_df.loc[
            mask,
            "血糖値"
        ]

        count_value = int(target.count())

        if count_value >= minimum_count:
            average_value = round(
                float(target.mean()),
                1
            )
        else:
            average_value = np.nan

        average_list.append(average_value)
        count_list.append(count_value)

    return average_list, count_list


def calculate_regression(x_values, y_values):
    """
    単回帰分析を行い、
    傾き、切片、相関係数、決定係数を返す。
    """
    x = np.asarray(x_values, dtype=float)
    y = np.asarray(y_values, dtype=float)

    if len(x) < 3:
        raise ValueError(
            "回帰分析には3件以上のデータが必要です。"
        )

    if np.allclose(x, x[0]):
        raise ValueError(
            "説明変数がすべて同じため、回帰分析できません。"
        )

    if np.allclose(y, y[0]):
        raise ValueError(
            "HA1cがすべて同じため、回帰分析できません。"
        )

    slope, intercept = np.polyfit(x, y, 1)
    predicted_y = slope * x + intercept

    correlation = float(
        np.corrcoef(x, y)[0, 1]
    )

    ss_res = float(
        np.sum((y - predicted_y) ** 2)
    )

    ss_tot = float(
        np.sum((y - np.mean(y)) ** 2)
    )

    r_squared = (
        np.nan
        if ss_tot == 0
        else 1 - ss_res / ss_tot
    )

    return {
        "slope": float(slope),
        "intercept": float(intercept),
        "correlation": correlation,
        "r_squared": float(r_squared),
        "predicted_y": predicted_y
    }


def replace_sheet(workbook, sheet_name, index=None):
    """
    指定シートを削除して同じ位置に作り直す。
    """
    if sheet_name in workbook.sheetnames:
        old_index = workbook.sheetnames.index(sheet_name)
        workbook.remove(workbook[sheet_name])

        if index is None:
            index = old_index

    if index is None:
        return workbook.create_sheet(sheet_name)

    return workbook.create_sheet(sheet_name, index)


def add_dataframe_to_sheet(
    worksheet,
    dataframe,
    start_row=1,
    start_col=1
):
    """
    DataFrameを指定位置からワークシートへ書き込む。
    """
    for row_offset, row in enumerate(
        dataframe_to_rows(
            dataframe,
            index=False,
            header=True
        ),
        start=start_row
    ):
        for col_offset, value in enumerate(
            row,
            start=start_col
        ):
            worksheet.cell(
                row=row_offset,
                column=col_offset,
                value=value
            )


def create_regression_graph(
    dataframe,
    average_column,
    average_days,
    image_path,
    title_prefix="HA1cと測定日前"
):
    """
    指定日数の平均血糖値について、
    散布図、回帰直線、回帰結果をPNG出力する。
    """
    plot_data = dataframe.dropna(
        subset=["HA1c", average_column]
    ).copy()

    x = plot_data[average_column].to_numpy(dtype=float)
    y = plot_data["HA1c"].to_numpy(dtype=float)

    result = calculate_regression(x, y)

    x_line = np.linspace(
        x.min(),
        x.max(),
        200
    )

    y_line = (
        result["slope"] * x_line
        + result["intercept"]
    )

    fig, ax = plt.subplots(
        figsize=(13.33, 7.5)
    )

    ax.scatter(
        x,
        y,
        s=80,
        label="HA1c測定値"
    )

    ax.plot(
        x_line,
        y_line,
        linewidth=2.5,
        label="回帰直線"
    )

    for _, row in plot_data.iterrows():
        ax.annotate(
            row["日付"].strftime("%Y-%m-%d"),
            (
                row[average_column],
                row["HA1c"]
            ),
            xytext=(5, 5),
            textcoords="offset points",
            fontsize=8
        )

    ax.set_title(
        f"{title_prefix}{average_days}日平均血糖値の関係",
        fontsize=20,
        pad=18
    )

    ax.set_xlabel(
        f"測定日前{average_days}日間の平均血糖値（mg/dL）",
        fontsize=13
    )

    ax.set_ylabel(
        "HA1c（%）",
        fontsize=13
    )

    ax.grid(
        True,
        linestyle="--",
        alpha=0.4
    )

    ax.legend(
        loc="upper left"
    )

    analysis_text = (
        f"回帰式：HA1c = "
        f"{result['slope']:.4f} × 平均血糖値 "
        f"{result['intercept']:+.3f}\n"
        f"回帰係数（傾き） = {result['slope']:.4f}\n"
        f"切片 = {result['intercept']:.3f}\n"
        f"決定係数 R² = {result['r_squared']:.3f}\n"
        f"相関係数 r = {result['correlation']:.3f}\n"
        f"データ数 n = {len(plot_data)}"
    )

    ax.text(
        0.98,
        0.04,
        analysis_text,
        transform=ax.transAxes,
        fontsize=12,
        verticalalignment="bottom",
        horizontalalignment="right",
        bbox={
            "boxstyle": "round",
            "facecolor": "white",
            "alpha": 0.85
        }
    )

    fig.tight_layout()

    fig.savefig(
        image_path,
        dpi=180,
        bbox_inches="tight"
    )

    plt.close(fig)

    return result, len(plot_data)


def add_graph_slide(
    presentation,
    title,
    image_path
):
    """
    タイトルとグラフ画像をPowerPointに追加する。
    """
    blank_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank_layout)

    title_box = slide.shapes.add_textbox(
        Inches(0.4),
        Inches(0.12),
        Inches(12.5),
        Inches(0.55)
    )

    title_frame = title_box.text_frame
    title_frame.clear()

    paragraph = title_frame.paragraphs[0]
    paragraph.text = title
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.size = Pt(24)
    paragraph.font.bold = True

    slide.shapes.add_picture(
        image_path,
        Inches(0.45),
        Inches(0.75),
        width=Inches(12.43),
        height=Inches(6.55)
    )


# ============================================================
# データ読み込み
# ============================================================

print("HA1cデータを読み込んでいます。")
df_hba1c = load_hba1c_data(hba1c_path)

print("血糖値データを読み込んでいます。")
df_glucose = load_glucose_data(glucose_path)

if df_hba1c.empty:
    raise ValueError(
        "HA1cデータがありません。"
    )

if df_glucose.empty:
    raise ValueError(
        "血糖値データがありません。"
    )

print(
    f"HA1cデータ期間: "
    f"{df_hba1c['日付'].min():%Y-%m-%d} ～ "
    f"{df_hba1c['日付'].max():%Y-%m-%d}"
)

print(
    f"血糖値データ期間: "
    f"{df_glucose['日付'].min():%Y-%m-%d} ～ "
    f"{df_glucose['日付'].max():%Y-%m-%d}"
)


# ============================================================
# 60日平均を計算
# ============================================================

average_60, count_60 = calculate_average_glucose(
    df_hba1c,
    df_glucose,
    STANDARD_DAYS
)

df_hba1c["前60日平均血糖値"] = average_60
df_hba1c["血糖値データ数"] = count_60


# ============================================================
# 30日～120日の平均血糖値を計算
# ============================================================

analysis_df = df_hba1c[
    ["日付", "HA1c"]
].copy()

for days in AVERAGE_DAYS_LIST:

    average_values, count_values = calculate_average_glucose(
        df_hba1c,
        df_glucose,
        days
    )

    analysis_df[
        f"平均血糖値_{days}日"
    ] = average_values

    analysis_df[
        f"データ数_{days}日"
    ] = count_values


# ============================================================
# 全日数で共通して使用できるHA1c測定日を抽出
# ============================================================

average_columns = [
    f"平均血糖値_{days}日"
    for days in AVERAGE_DAYS_LIST
]

common_analysis_df = analysis_df.dropna(
    subset=["HA1c"] + average_columns
).copy()

if len(common_analysis_df) < 3:
    raise ValueError(
        "30日～120日の全期間で共通して使用できる"
        "HA1cデータが3件未満です。\n"
        "MIN_COVERAGE_RATIOを確認してください。"
    )

print(
    "平均日数比較に使用する共通データ数:",
    len(common_analysis_df)
)

print(
    "共通データ期間:",
    common_analysis_df["日付"].min().strftime("%Y-%m-%d"),
    "～",
    common_analysis_df["日付"].max().strftime("%Y-%m-%d")
)


# ============================================================
# 各平均日数の回帰分析
# ============================================================

comparison_rows = []

for days in AVERAGE_DAYS_LIST:

    average_column = f"平均血糖値_{days}日"

    x = common_analysis_df[
        average_column
    ].to_numpy(dtype=float)

    y = common_analysis_df[
        "HA1c"
    ].to_numpy(dtype=float)

    result = calculate_regression(x, y)

    comparison_rows.append({
        "平均日数": days,
        "相関係数r": result["correlation"],
        "決定係数R2": result["r_squared"],
        "回帰係数": result["slope"],
        "切片": result["intercept"],
        "データ数": len(common_analysis_df)
    })


comparison_df = pd.DataFrame(comparison_rows)

# 決定係数が高い順に順位付け
comparison_df["順位"] = (
    comparison_df["決定係数R2"]
    .rank(
        method="min",
        ascending=False
    )
    .astype(int)
)

best_row = comparison_df.loc[
    comparison_df["決定係数R2"].idxmax()
]

best_days = int(best_row["平均日数"])
best_r = float(best_row["相関係数r"])
best_r2 = float(best_row["決定係数R2"])
best_slope = float(best_row["回帰係数"])
best_intercept = float(best_row["切片"])

best_average_column = (
    f"平均血糖値_{best_days}日"
)

best_regression_png = os.path.join(
    output_dir,
    f"HA1c_{best_days}日平均血糖値_回帰分析.png"
)

print("=" * 60)
print(f"最適平均日数: {best_days}日")
print(f"相関係数 r: {best_r:.3f}")
print(f"決定係数 R²: {best_r2:.3f}")
print(f"回帰係数（傾き）: {best_slope:.4f}")
print(f"切片: {best_intercept:.3f}")
print(f"共通データ数: {len(common_analysis_df)}")
print("=" * 60)


# ============================================================
# ExcelのSheet2と平均日数分析シートを更新
# ============================================================

wb = openpyxl.load_workbook(hba1c_path)

sheet2_index = (
    wb.sheetnames.index("Sheet2")
    if "Sheet2" in wb.sheetnames
    else None
)

ws = replace_sheet(
    wb,
    "Sheet2",
    sheet2_index
)

df_excel = df_hba1c.copy()
df_excel["日付"] = df_excel["日付"].dt.date

add_dataframe_to_sheet(
    ws,
    df_excel
)

ws.column_dimensions["A"].width = 14
ws.column_dimensions["B"].width = 12
ws.column_dimensions["C"].width = 20
ws.column_dimensions["D"].width = 18


# ============================================================
# Excelに60日平均の2軸折れ線グラフを作成
# ============================================================

graph_obj1 = LineChart()
graph_obj1.title = "HA1cと平均血糖値（測定日前60日間）"
graph_obj1.x_axis.title = "日付"
graph_obj1.y_axis.title = "HA1c（%）"

graph_obj2 = LineChart()
graph_obj2.y_axis.title = "60日平均血糖値（mg/dL）"

v1 = Reference(
    ws,
    min_col=2,
    min_row=1,
    max_row=ws.max_row
)

categories = Reference(
    ws,
    min_col=1,
    min_row=2,
    max_row=ws.max_row
)

v2 = Reference(
    ws,
    min_col=3,
    min_row=1,
    max_row=ws.max_row
)

graph_obj1.add_data(
    v1,
    titles_from_data=True
)

graph_obj1.set_categories(categories)

graph_obj2.add_data(
    v2,
    titles_from_data=True
)

graph_obj2.set_categories(categories)

for series in graph_obj1.series:
    series.smooth = False

for series in graph_obj2.series:
    series.smooth = False

# HA1c（左側Y軸）の表示範囲
graph_obj1.y_axis.scaling.min = 4
graph_obj1.y_axis.scaling.max = 9

# 60日平均血糖（右側Y軸）の表示範囲
graph_obj2.y_axis.axId = 200
graph_obj2.y_axis.crosses = "max"
graph_obj2.y_axis.majorGridlines = None
graph_obj2.y_axis.scaling.min = 80
graph_obj2.y_axis.scaling.max = 160

graph_obj1.width = 25
graph_obj1.height = 15

graph_obj1 += graph_obj2
ws.add_chart(graph_obj1, "F2")


# ============================================================
# Excelに平均日数分析結果を保存
# ============================================================

analysis_ws = replace_sheet(
    wb,
    ANALYSIS_SHEET_NAME
)

excel_comparison_df = comparison_df.copy()

excel_comparison_df["相関係数r"] = (
    excel_comparison_df["相関係数r"].round(4)
)

excel_comparison_df["決定係数R2"] = (
    excel_comparison_df["決定係数R2"].round(4)
)

excel_comparison_df["回帰係数"] = (
    excel_comparison_df["回帰係数"].round(5)
)

excel_comparison_df["切片"] = (
    excel_comparison_df["切片"].round(4)
)

add_dataframe_to_sheet(
    analysis_ws,
    excel_comparison_df,
    start_row=1,
    start_col=1
)

analysis_ws["I1"] = "最適平均日数"
analysis_ws["J1"] = best_days

analysis_ws["I2"] = "相関係数r"
analysis_ws["J2"] = round(best_r, 4)

analysis_ws["I3"] = "決定係数R2"
analysis_ws["J3"] = round(best_r2, 4)

analysis_ws["I4"] = "回帰係数"
analysis_ws["J4"] = round(best_slope, 5)

analysis_ws["I5"] = "切片"
analysis_ws["J5"] = round(best_intercept, 4)

analysis_ws["I6"] = "共通データ数"
analysis_ws["J6"] = len(common_analysis_df)

analysis_ws["I7"] = "共通期間開始"
analysis_ws["J7"] = (
    common_analysis_df["日付"].min().date()
)

analysis_ws["I8"] = "共通期間終了"
analysis_ws["J8"] = (
    common_analysis_df["日付"].max().date()
)

analysis_ws["I9"] = "必要データ割合"
analysis_ws["J9"] = MIN_COVERAGE_RATIO

analysis_ws.column_dimensions["A"].width = 12
analysis_ws.column_dimensions["B"].width = 14
analysis_ws.column_dimensions["C"].width = 14
analysis_ws.column_dimensions["D"].width = 14
analysis_ws.column_dimensions["E"].width = 12
analysis_ws.column_dimensions["F"].width = 12
analysis_ws.column_dimensions["G"].width = 10
analysis_ws.column_dimensions["I"].width = 18
analysis_ws.column_dimensions["J"].width = 16

wb.active = wb.sheetnames.index("Sheet2")
wb.save(hba1c_path)

print(f"Excel保存: {hba1c_path}")

# ============================================================
# データ表1.xlsx Sheet5 に日別HA1c予測値を保存
# ============================================================

PREDICTION_START_DATE = pd.Timestamp("2025-01-01")

prediction_rows = []

# 回帰分析で最適と判定された日数を使用
prediction_days = best_days
prediction_slope = best_slope
prediction_intercept = best_intercept

minimum_count = 21

for _, row in df_glucose.iterrows():

    target_date = row["日付"]

    if target_date < PREDICTION_START_DATE:
        continue

    # 回帰分析と同じ条件：
    # 当日は含めず、直前 prediction_days 日間を使用
    start_date = target_date - timedelta(days=prediction_days)

    mask = (
        (df_glucose["日付"] >= start_date)
        & (df_glucose["日付"] < target_date)
    )

    target_glucose = df_glucose.loc[
        mask,
        "血糖値"
    ]

    count_value = int(target_glucose.count())

    if count_value >= minimum_count:

        average_glucose = float(
            target_glucose.mean()
        )

        predicted_hba1c = (
            prediction_slope * average_glucose
            + prediction_intercept
        )

        predicted_hba1c = round(
            predicted_hba1c,
            2
        )

    else:
        predicted_hba1c = np.nan

    prediction_rows.append({
        "日付": target_date.date(),
        "血糖値": row["血糖値"],
        "HA1c予測値": predicted_hba1c
    })


prediction_df = pd.DataFrame(
    prediction_rows
)


# ============================================================
# Sheet5を書き換える
# ============================================================

wb_glucose = openpyxl.load_workbook(
    glucose_path
)

if "Sheet5" in wb_glucose.sheetnames:

    sheet5_index = wb_glucose.sheetnames.index(
        "Sheet5"
    )

    wb_glucose.remove(
        wb_glucose["Sheet5"]
    )

    ws5 = wb_glucose.create_sheet(
        "Sheet5",
        sheet5_index
    )

else:
    ws5 = wb_glucose.create_sheet(
        "Sheet5"
    )


add_dataframe_to_sheet(
    ws5,
    prediction_df
)

ws5.column_dimensions["A"].width = 14
ws5.column_dimensions["B"].width = 12
ws5.column_dimensions["C"].width = 16

# 日付表示
for cell in ws5["A"][1:]:
    cell.number_format = "yyyy/mm/dd"

# HA1c予測値を小数第2位まで表示
for cell in ws5["C"][1:]:
    cell.number_format = "0.00"


wb_glucose.save(
    glucose_path
)

print("=" * 60)
print("日別HA1c予測値を作成しました。")
print(f"予測開始日: {PREDICTION_START_DATE:%Y-%m-%d}")
print(f"使用平均日数: {prediction_days}日")
print(f"回帰係数: {prediction_slope:.4f}")
print(f"切片: {prediction_intercept:.3f}")
print(f"出力件数: {len(prediction_df)}")
print(f"保存先: {glucose_path} / Sheet5")
print("=" * 60)

# ============================================================
# グラフ：血糖値から予測されるHA1c
# ============================================================

plot_prediction_df = prediction_df.dropna(
    subset=["血糖値", "HA1c予測値"]
).copy()

if plot_prediction_df.empty:
    raise ValueError(
        "HA1c予測グラフを作成できるデータがありません。"
    )

fig, ax1 = plt.subplots(
    figsize=(13.33, 7.5)
)

ax2 = ax1.twinx()

# 実測血糖値
line1 = ax1.plot(
    plot_prediction_df["日付"],
    plot_prediction_df["血糖値"],
    linewidth=1.8,
    label="血糖値"
)

# HA1c予測値
line2 = ax2.plot(
    plot_prediction_df["日付"],
    plot_prediction_df["HA1c予測値"],
    color="tab:orange",
    linewidth=3,
    label="HA1c予測値"
)

ax1.set_title(
    "血糖値から予測されるHA1c",
    fontsize=20,
    pad=18
)

ax1.set_xlabel(
    "日付",
    fontsize=13
)

ax1.set_ylabel(
    "朝食前血糖値（mg/dL）",
    fontsize=13
)

ax2.set_ylabel(
    "HA1c予測値（%）",
    fontsize=13
)

# 血糖値軸
ax1.set_ylim(
    0,
    max(
        250,
        plot_prediction_df["血糖値"].max() + 20
    )
)

# HA1c軸
ax2.set_ylim(
    5.8,
    8.0
)

ax1.xaxis.set_major_locator(
    mdates.MonthLocator()
)

ax1.xaxis.set_major_formatter(
    mdates.DateFormatter("%Y/%m/%d")
)

plt.setp(
    ax1.get_xticklabels(),
    rotation=45,
    ha="right"
)

ax1.grid(
    True,
    linestyle="--",
    alpha=0.4
)

lines = line1 + line2

labels = [
    line.get_label()
    for line in lines
]

ax1.legend(
    lines,
    labels,
    loc="upper left"
)

fig.tight_layout()

fig.savefig(
    prediction_png,
    dpi=180,
    bbox_inches="tight"
)

plt.close(fig)

print(
    f"HA1c予測グラフPNG保存: {prediction_png}"
)

# ============================================================
# グラフ1：HA1cと60日平均血糖値の時系列
# ============================================================

plot_60_df = df_hba1c.dropna(
    subset=["HA1c", "前60日平均血糖値"]
).copy()

if plot_60_df.empty:
    raise ValueError(
        "HA1cと前60日平均血糖値の両方が存在する"
        "データがありません。"
    )

fig, ax1 = plt.subplots(
    figsize=(13.33, 7.5)
)

ax2 = ax1.twinx()

line1 = ax1.plot(
    plot_60_df["日付"],
    plot_60_df["HA1c"],
    color="red",
    marker="o",
    linewidth=3,
    label="HA1c"
)

line2 = ax2.plot(
    plot_60_df["日付"],
    plot_60_df["前60日平均血糖値"],
    color="blue",
    marker="s",
    linewidth=3,
    label="前60日平均血糖値"
)

ax1.set_title(
    "HA1cと測定日前60日間の平均血糖値",
    fontsize=20,
    pad=18
)

ax1.set_xlabel(
    "HA1c測定日",
    fontsize=13
)

ax1.set_ylabel(
    "HA1c（%）",
    fontsize=13
)

ax2.set_ylabel(
    "60日平均血糖値（mg/dL）",
    fontsize=13
)

hba1c_min = plot_60_df["HA1c"].min()
hba1c_max = plot_60_df["HA1c"].max()

ax1.set_ylim(
    max(0, hba1c_min - 0.5),
    hba1c_max + 0.5
)

glucose_min = (
    plot_60_df["前60日平均血糖値"].min()
)

glucose_max = (
    plot_60_df["前60日平均血糖値"].max()
)

margin = max(
    10,
    (glucose_max - glucose_min) * 0.15
)

ax2.set_ylim(
    max(0, glucose_min - margin),
    glucose_max + margin
)

ax1.xaxis.set_major_locator(
    mdates.AutoDateLocator()
)

ax1.xaxis.set_major_formatter(
    mdates.DateFormatter("%Y-%m-%d")
)

plt.setp(
    ax1.get_xticklabels(),
    rotation=45,
    ha="right"
)

ax1.grid(
    True,
    linestyle="--",
    alpha=0.4
)

lines = line1 + line2
labels = [
    line.get_label()
    for line in lines
]

ax1.legend(
    lines,
    labels,
    loc="upper left"
)

fig.tight_layout()

fig.savefig(
    time_series_png,
    dpi=180,
    bbox_inches="tight"
)

plt.close(fig)

print(f"時系列PNG保存: {time_series_png}")


# ============================================================
# グラフ2：60日平均の散布図・回帰分析
# ============================================================

result_60, count_result_60 = create_regression_graph(
    dataframe=plot_60_df,
    average_column="前60日平均血糖値",
    average_days=STANDARD_DAYS,
    image_path=regression_60_png
)

print(
    f"60日平均回帰分析PNG保存: {regression_60_png}"
)


# ============================================================
# グラフ3：平均日数別の相関係数
# ============================================================

fig, ax = plt.subplots(
    figsize=(13.33, 7.5)
)

ax.plot(
    comparison_df["平均日数"],
    comparison_df["相関係数r"],
    color="red",
    marker="o",
    linewidth=3,
    label="相関係数 r"
)

ax.axvline(
    best_days,
    color="green",
    linestyle="--",
    linewidth=2,
    label=f"最適日数：{best_days}日"
)

ax.scatter(
    [best_days],
    [best_r],
    color="green",
    s=150,
    zorder=5
)

ax.annotate(
    (
        f"最適日数：{best_days}日\n"
        f"r = {best_r:.3f}"
    ),
    xy=(best_days, best_r),
    xytext=(15, -10),
    textcoords="offset points",
    fontsize=12,
    bbox={
        "boxstyle": "round",
        "facecolor": "white",
        "alpha": 0.9
    }
)

ax.set_title(
    "平均日数別の相関係数",
    fontsize=20,
    pad=18
)

ax.set_xlabel(
    "HA1c測定日前の平均日数",
    fontsize=13
)

ax.set_ylabel(
    "相関係数 r",
    fontsize=13
)

ax.set_xticks(AVERAGE_DAYS_LIST)
ax.set_ylim(0, 1)

ax.grid(
    True,
    linestyle="--",
    alpha=0.4
)

ax.legend(
    loc="best"
)

fig.tight_layout()

fig.savefig(
    correlation_comparison_png,
    dpi=180,
    bbox_inches="tight"
)

plt.close(fig)

print(
    f"相関係数比較PNG保存: "
    f"{correlation_comparison_png}"
)


# ============================================================
# グラフ4：平均日数別の決定係数
# ============================================================

fig, ax = plt.subplots(
    figsize=(13.33, 7.5)
)

ax.plot(
    comparison_df["平均日数"],
    comparison_df["決定係数R2"],
    color="blue",
    marker="s",
    linewidth=3,
    label="決定係数 R²"
)

ax.axvline(
    best_days,
    color="green",
    linestyle="--",
    linewidth=2,
    label=f"最適日数：{best_days}日"
)

ax.scatter(
    [best_days],
    [best_r2],
    color="green",
    s=150,
    zorder=5
)

ax.annotate(
    (
        f"最適日数：{best_days}日\n"
        f"R² = {best_r2:.3f}\n"
        f"傾き = {best_slope:.4f}"
    ),
    xy=(best_days, best_r2),
    xytext=(15, -10),
    textcoords="offset points",
    fontsize=12,
    bbox={
        "boxstyle": "round",
        "facecolor": "white",
        "alpha": 0.9
    }
)

ax.set_title(
    "平均日数別の決定係数",
    fontsize=20,
    pad=18
)

ax.set_xlabel(
    "HA1c測定日前の平均日数",
    fontsize=13
)

ax.set_ylabel(
    "決定係数 R²",
    fontsize=13
)

ax.set_xticks(AVERAGE_DAYS_LIST)
ax.set_ylim(0, 1)

ax.grid(
    True,
    linestyle="--",
    alpha=0.4
)

ax.legend(
    loc="best"
)

fig.tight_layout()

fig.savefig(
    r2_comparison_png,
    dpi=180,
    bbox_inches="tight"
)

plt.close(fig)

print(
    f"決定係数比較PNG保存: "
    f"{r2_comparison_png}"
)


# ============================================================
# グラフ5：最適日数の散布図・回帰分析
# ============================================================

best_result, best_count = create_regression_graph(
    dataframe=common_analysis_df,
    average_column=best_average_column,
    average_days=best_days,
    image_path=best_regression_png,
    title_prefix="HA1cと最適期間・測定日前"
)

print(
    f"最適日数回帰分析PNG保存: "
    f"{best_regression_png}"
)


# ============================================================
# PowerPoint作成
# ============================================================

prs = Presentation()

prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 1枚目
add_graph_slide(
    prs,
    "HA1cと測定日前60日間の平均血糖値",
    time_series_png
)

# 2枚目
add_graph_slide(
    prs,
    "HA1cと60日平均血糖値の散布図・回帰分析",
    regression_60_png
)

# 3枚目
add_graph_slide(
    prs,
    "平均日数別の相関係数",
    correlation_comparison_png
)

# 4枚目
add_graph_slide(
    prs,
    "平均日数別の決定係数",
    r2_comparison_png
)

# 5枚目
add_graph_slide(
    prs,
    f"最適平均日数（{best_days}日）の散布図・回帰分析",
    best_regression_png
)

# 6枚目
add_graph_slide(
    prs,
    "血糖値から予測されるHA1c",
    prediction_png
)

prs.save(pptx_path)

print(f"PowerPoint保存: {pptx_path}")


# ============================================================
# PowerPointを開く
# ============================================================

os.startfile(pptx_path)

print("処理が完了しました。")


# label="HA1c予測値"
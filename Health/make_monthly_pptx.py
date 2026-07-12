import os
import traceback
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


# =========================
# 月次グラフの定義
# =========================

ITEMS = {
    1: {
        "title": "体重（月次平均値 ±1σ）",
        "filename": "1_体重_月次平均.png",
    },
    2: {
        "title": "血糖値（月次平均値 ±1σ）",
        "filename": "2_血糖値_月次平均.png",
    },
    3: {
        "title": "血圧収縮期（月次平均値 ±1σ）",
        "filename": "3_血圧収縮期_月次平均.png",
    },
    4: {
        "title": "血圧拡張期（月次平均値 ±1σ）",
        "filename": "4_血圧拡張期_月次平均.png",
    },
    5: {
        "title": "心拍数（月次平均値 ±1σ）",
        "filename": "5_心拍数_月次平均.png",
    },
    6: {
        "title": "中程度運動量（月次平均値 ±1σ）",
        "filename": "6_中程度運動量_月次平均.png",
    },
    7: {
        "title": "運動消費エネルギー（月次平均値 ±1σ）",
        "filename": "7_運動消費エネルギー_月次平均.png",
    },
    8: {
        "title": "歩数（月次平均値 ±1σ）",
        "filename": "8_歩数_月次平均.png",
    },
}


# =========================
# パス設定
# =========================

def get_onedrive_path(*parts):
    return Path(os.environ["OneDrive"]).joinpath(*parts)


BASE_HEALTH = get_onedrive_path(
    "ドキュメント",
    "PythonWork",
    "Health",
)

PNG_DIR = BASE_HEALTH / "monthly_graph_png"

PPTX_PATH = BASE_HEALTH / "健康管理_月次推移.pptx"

ERROR_LOG_PATH = BASE_HEALTH / "make_monthly_pptx_error.txt"


# =========================
# スライド作成
# =========================

def add_slide(prs, title, png_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # スライド上部のタイトル
    title_box = slide.shapes.add_textbox(
        Inches(0.4),
        Inches(0.15),
        Inches(9.2),
        Inches(0.55),
    )

    text_frame = title_box.text_frame
    text_frame.clear()

    paragraph = text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.size = Pt(24)
    paragraph.font.bold = True

    # グラフ画像を貼り付け
    slide.shapes.add_picture(
        str(png_path),
        Inches(0.35),
        Inches(0.75),
        width=Inches(9.3),
    )


# =========================
# PowerPoint作成
# =========================

def make_pptx():
    if not PNG_DIR.exists():
        raise FileNotFoundError(
            f"月次グラフフォルダーがありません: {PNG_DIR}"
        )

    prs = Presentation()

    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for item_number, item_info in ITEMS.items():
        png_path = PNG_DIR / item_info["filename"]

        if not png_path.exists():
            raise FileNotFoundError(
                f"PNGファイルがありません: {png_path}"
            )

        print(
            f"スライド作成: {item_number} "
            f"{item_info['title']}"
        )

        add_slide(
            prs=prs,
            title=item_info["title"],
            png_path=png_path,
        )

    prs.save(PPTX_PATH)

    print("PowerPoint作成完了:", PPTX_PATH)

    return PPTX_PATH


# =========================
# メイン処理
# =========================

def main():
    try:
        pptx_path = make_pptx()

        # 作成したPowerPointを自動で開く
        os.startfile(pptx_path)

    except Exception:
        with open(
            ERROR_LOG_PATH,
            "w",
            encoding="utf-8",
        ) as log_file:
            traceback.print_exc(file=log_file)

        print("エラー内容を保存しました:", ERROR_LOG_PATH)
        raise


if __name__ == "__main__":
    main()
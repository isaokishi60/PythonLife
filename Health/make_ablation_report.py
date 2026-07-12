import os
import sys
import time
import argparse
import subprocess
import traceback
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

try:
    import win32com.client
except ImportError:
    raise ImportError(
        "pywin32 が必要です。pip install pywin32 を実行してください。"
    )


ITEMS = {
    1: "体重",
    2: "血糖値",
    3: "血圧・心拍数",
    4: "中程度運動量",
    5: "運動消費エネルギー",
    6: "歩数",
}

# Sheet4 内の列名に多少の違いがあっても認識できるようにしています。
COLUMN_ALIASES = {
    "日付": ["日付", "年月日", "Date"],
    "体重": ["体重", "体重(kg)", "体重（kg）"],
    "血糖値": ["血糖値", "血糖", "朝食前血糖", "空腹時血糖"],
    "血圧収縮期": ["血圧収縮期", "収縮期血圧", "最高血圧", "SBP"],
    "血圧拡張期": ["血圧拡張期", "拡張期血圧", "最低血圧", "DBP"],
    "心拍数": ["心拍数", "脈拍", "HR"],
    "中程度運動量": [
        "中程度運動量（分）",
        "中程度運動量(分)",
        "中程度運動量",
        "中強度運動量（分）",
    ],
    "運動消費エネルギー": [
        "運動消費kcal",
        "運動消費Kcal",
        "運動消費エネルギー",
        "運動消費カロリー",
    ],
    "歩数": ["歩数", "Steps"],
}

SUMMARY_METRICS = [
    ("体重", "kg", 2),
    ("血糖値", "mg/dL", 1),
    ("血圧収縮期", "mmHg", 1),
    ("血圧拡張期", "mmHg", 1),
    ("心拍数", "回/分", 1),
    ("中程度運動量", "分/日", 1),
    ("運動消費エネルギー", "kcal/日", 1),
    ("歩数", "歩/日", 0),
]


def get_onedrive_path(*parts):
    return Path(os.environ["OneDrive"]).joinpath(*parts)


BASE_HEALTH = get_onedrive_path("ドキュメント", "PythonWork", "Health")
PNG_DIR = BASE_HEALTH / "health_graph_png"

NORMAL_PPTX_PATH = BASE_HEALTH / "健康管理グラフ.pptx"
ABLATION_PPTX_PATH = BASE_HEALTH / "アブレーション前後比較レポート.pptx"

EXCEL_PATH = get_onedrive_path(
    "ドキュメント", "PythonWork", "ExcelDATA", "データ表1.xlsx"
)

GRAPH_SCRIPT = (
    Path(__file__).parent
    / "04_charts"
    / "Graph_from_Excel(Complete_2025_11_03).py"
)


def normalize_text(value):
    return str(value).strip().replace(" ", "").replace("　", "")


def find_column(columns, logical_name):
    normalized_columns = {
        normalize_text(column): column for column in columns
    }

    for alias in COLUMN_ALIASES[logical_name]:
        key = normalize_text(alias)
        if key in normalized_columns:
            return normalized_columns[key]

    return None


def run_graph_script(start_date, end_date, item):
    cmd = [
        sys.executable,
        str(GRAPH_SCRIPT),
        "--start-date",
        start_date,
        "--end-date",
        end_date,
        "--item",
        str(item),
    ]

    print("グラフ作成:", ITEMS[item])
    result = subprocess.run(
        cmd,
        text=True,
        capture_output=True,
        encoding="cp932",
        errors="replace",
    )

    if result.stdout:
        print(result.stdout)

    if result.returncode != 0:
        if result.stderr:
            print(result.stderr)
        raise RuntimeError(
            f"グラフ作成に失敗しました: item={item}"
        )


def export_excel_chart_to_png(item):
    PNG_DIR.mkdir(parents=True, exist_ok=True)
    png_path = PNG_DIR / f"{item}_{ITEMS[item]}.png"

    # 前回のPNGを誤って再利用しないように削除します。
    if png_path.exists():
        png_path.unlink()

    excel = win32com.client.DispatchEx("Excel.Application")
    excel.Visible = False
    excel.DisplayAlerts = False
    excel.ScreenUpdating = True

    wb = None

    try:
        wb = excel.Workbooks.Open(
            str(EXCEL_PATH),
            UpdateLinks=0,
            ReadOnly=True,
        )
        time.sleep(5)

        ws = wb.Worksheets("Sheet3")
        ws.Activate()

        charts = ws.ChartObjects()
        print("ChartObjects Count:", charts.Count)

        if charts.Count == 0:
            raise RuntimeError("Sheet3 にグラフがありません")

        time.sleep(5)

        for retry in range(10):
            try:
                excel.CalculateFullRebuild()
                excel.CalculateUntilAsyncQueriesDone()
                time.sleep(1)

                chart_obj = ws.ChartObjects(1)
                chart_obj.Activate()

                chart = chart_obj.Chart
                chart.Export(str(png_path))

                if png_path.exists() and png_path.stat().st_size > 0:
                    break

            except Exception as exc:
                print(
                    f"グラフPNG出力リトライ "
                    f"{retry + 1}/10: {repr(exc)}"
                )
                time.sleep(5)
        else:
            raise RuntimeError(
                f"グラフPNG出力に失敗しました: item={item}"
            )

    finally:
        if wb is not None:
            wb.Close(SaveChanges=False)
        excel.Quit()

    print("PNG保存:", png_path)
    return png_path


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(0.22),
        Inches(9.0),
        Inches(0.55),
    )
    paragraph = title_box.text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.name = "Yu Gothic"
    paragraph.font.size = Pt(25)
    paragraph.font.bold = True

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.52),
            Inches(0.73),
            Inches(8.9),
            Inches(0.3),
        )
        paragraph = subtitle_box.text_frame.paragraphs[0]
        paragraph.text = subtitle
        paragraph.font.name = "Yu Gothic"
        paragraph.font.size = Pt(10)


def add_footer(slide):
    footer_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(7.15),
        Inches(9.0),
        Inches(0.2),
    )
    paragraph = footer_box.text_frame.paragraphs[0]
    paragraph.text = (
        "診察時の説明用資料です。医学的評価は担当医にご確認ください。"
    )
    paragraph.font.name = "Yu Gothic"
    paragraph.font.size = Pt(8)
    paragraph.alignment = PP_ALIGN.RIGHT


def add_graph_slide(prs, title, png_path, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title, subtitle)

    slide.shapes.add_picture(
        str(png_path),
        Inches(0.5),
        Inches(1.0),
        width=Inches(9.0),
    )

    if subtitle:
        add_footer(slide)


def load_ablation_data(start_date, end_date):
    print("比較データ読込:", EXCEL_PATH)

    df = pd.read_excel(EXCEL_PATH, sheet_name="Sheet4")

    date_column = find_column(df.columns, "日付")
    if date_column is None:
        raise RuntimeError(
            "Sheet4 に日付列が見つかりません。"
            f" 列名: {list(df.columns)}"
        )

    rename_map = {date_column: "日付"}

    for logical_name, _, _ in SUMMARY_METRICS:
        actual_column = find_column(df.columns, logical_name)
        if actual_column is not None:
            rename_map[actual_column] = logical_name

    df = df.rename(columns=rename_map).copy()
    df["日付"] = pd.to_datetime(df["日付"], errors="coerce")
    df = df.dropna(subset=["日付"])

    start_timestamp = pd.Timestamp(start_date)
    end_timestamp = pd.Timestamp(end_date)

    df = df.loc[
        (df["日付"] >= start_timestamp)
        & (df["日付"] <= end_timestamp)
    ].copy()

    for logical_name, _, _ in SUMMARY_METRICS:
        if logical_name in df.columns:
            df[logical_name] = pd.to_numeric(
                df[logical_name],
                errors="coerce",
            )

    if df.empty:
        raise RuntimeError(
            f"Sheet4 に指定期間のデータがありません: "
            f"{start_date} ～ {end_date}"
        )

    print("比較対象行数:", len(df))
    return df


def calculate_ablation_summary(df, ablation_date):
    ablation_timestamp = pd.Timestamp(ablation_date)

    before = df.loc[df["日付"] < ablation_timestamp]
    after = df.loc[df["日付"] >= ablation_timestamp]

    if before.empty:
        raise RuntimeError(
            "アブレーション前のデータがありません。"
        )

    if after.empty:
        raise RuntimeError(
            "アブレーション後のデータがありません。"
        )

    summary_rows = []

    for metric_name, unit, decimals in SUMMARY_METRICS:
        if metric_name not in df.columns:
            print("比較対象外（列なし）:", metric_name)
            continue

        before_values = before[metric_name].dropna()
        after_values = after[metric_name].dropna()

        before_mean = (
            float(before_values.mean())
            if not before_values.empty
            else None
        )
        after_mean = (
            float(after_values.mean())
            if not after_values.empty
            else None
        )

        if before_mean is not None and after_mean is not None:
            change = after_mean - before_mean
            change_percent = (
                change / before_mean * 100
                if before_mean != 0
                else None
            )
        else:
            change = None
            change_percent = None

        summary_rows.append(
            {
                "metric": metric_name,
                "unit": unit,
                "decimals": decimals,
                "before_mean": before_mean,
                "after_mean": after_mean,
                "change": change,
                "change_percent": change_percent,
                "before_n": int(before_values.count()),
                "after_n": int(after_values.count()),
            }
        )

    return summary_rows


def format_number(value, decimals):
    if value is None or pd.isna(value):
        return "―"

    if decimals == 0:
        return f"{value:,.0f}"

    return f"{value:,.{decimals}f}"


def format_change(row):
    if row["change"] is None:
        return "―"

    change_text = format_number(
        row["change"],
        row["decimals"],
    )

    if row["change"] > 0:
        change_text = "+" + change_text

    if row["change_percent"] is not None:
        return (
            f"{change_text} {row['unit']} "
            f"({row['change_percent']:+.1f}%)"
        )

    return f"{change_text} {row['unit']}"


def add_dashboard_slide(
    prs,
    summary_rows,
    start_date,
    end_date,
    ablation_date,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(
        slide,
        "アブレーション前後比較ダッシュボード",
        (
            f"観察期間: {start_date} ～ {end_date}　"
            f"アブレーション: {ablation_date}"
        ),
    )

    summary_map = {
        row["metric"]: row for row in summary_rows
    }

    card_metrics = [
        "体重",
        "血糖値",
        "血圧収縮期",
        "心拍数",
    ]

    card_lefts = [0.5, 2.8, 5.1, 7.4]

    for metric_name, left in zip(card_metrics, card_lefts):
        if metric_name not in summary_map:
            continue

        row = summary_map[metric_name]

        card = slide.shapes.add_textbox(
            Inches(left),
            Inches(1.1),
            Inches(2.1),
            Inches(1.55),
        )
        text_frame = card.text_frame
        text_frame.clear()

        paragraph = text_frame.paragraphs[0]
        paragraph.text = metric_name
        paragraph.font.name = "Yu Gothic"
        paragraph.font.size = Pt(15)
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER

        paragraph = text_frame.add_paragraph()
        paragraph.text = (
            "前 "
            f"{format_number(row['before_mean'], row['decimals'])}"
            f" {row['unit']}"
        )
        paragraph.font.name = "Yu Gothic"
        paragraph.font.size = Pt(12)
        paragraph.alignment = PP_ALIGN.CENTER

        paragraph = text_frame.add_paragraph()
        paragraph.text = (
            "後 "
            f"{format_number(row['after_mean'], row['decimals'])}"
            f" {row['unit']}"
        )
        paragraph.font.name = "Yu Gothic"
        paragraph.font.size = Pt(12)
        paragraph.alignment = PP_ALIGN.CENTER

        paragraph = text_frame.add_paragraph()
        paragraph.text = "変化 " + format_change(row)
        paragraph.font.name = "Yu Gothic"
        paragraph.font.size = Pt(10)
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER

    rows = len(summary_rows) + 1
    columns = 6

    table = slide.shapes.add_table(
        rows,
        columns,
        Inches(0.45),
        Inches(2.85),
        Inches(9.1),
        Inches(3.9),
    ).table

    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(1.55)
    table.columns[2].width = Inches(1.55)
    table.columns[3].width = Inches(2.2)
    table.columns[4].width = Inches(0.9)
    table.columns[5].width = Inches(0.9)

    headers = [
        "項目",
        "前平均",
        "後平均",
        "変化",
        "前n",
        "後n",
    ]

    for column_index, header in enumerate(headers):
        cell = table.cell(0, column_index)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.name = "Yu Gothic"
            paragraph.font.size = Pt(9)
            paragraph.font.bold = True
            paragraph.alignment = PP_ALIGN.CENTER

    for row_index, row in enumerate(summary_rows, start=1):
        values = [
            row["metric"],
            (
                f"{format_number(row['before_mean'], row['decimals'])}"
                f" {row['unit']}"
            ),
            (
                f"{format_number(row['after_mean'], row['decimals'])}"
                f" {row['unit']}"
            ),
            format_change(row),
            str(row["before_n"]),
            str(row["after_n"]),
        ]

        for column_index, value in enumerate(values):
            cell = table.cell(row_index, column_index)
            cell.text = value

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = "Yu Gothic"
                paragraph.font.size = Pt(8)
                paragraph.alignment = (
                    PP_ALIGN.LEFT
                    if column_index == 0
                    else PP_ALIGN.CENTER
                )

    add_footer(slide)


def add_overview_slide(
    prs,
    start_date,
    end_date,
    ablation_date,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(
        slide,
        "アブレーション前後比較レポート",
        f"観察期間: {start_date} ～ {end_date}",
    )

    left_box = slide.shapes.add_textbox(
        Inches(0.75),
        Inches(1.25),
        Inches(5.2),
        Inches(4.8),
    )
    text_frame = left_box.text_frame
    text_frame.clear()

    lines = [
        f"アブレーション実施日：{ablation_date}",
        "",
        "主要イベント",
        "・2026/02/16　人間ドック",
        "・2026/03/19　カルディオバージョン",
        "・2026/05/22　アブレーション",
        "",
        "比較方法",
        "・前：アブレーション日の前日まで",
        "・後：アブレーション当日以降",
        "・指定期間内の日々の測定値を比較",
    ]

    for index, line in enumerate(lines):
        paragraph = (
            text_frame.paragraphs[0]
            if index == 0
            else text_frame.add_paragraph()
        )
        paragraph.text = line
        paragraph.font.name = "Yu Gothic"
        paragraph.font.size = Pt(15 if index == 0 else 13)

        if line in {"主要イベント", "比較方法"}:
            paragraph.font.bold = True
            paragraph.space_before = Pt(10)

    right_box = slide.shapes.add_textbox(
        Inches(6.2),
        Inches(1.35),
        Inches(3.1),
        Inches(4.5),
    )
    text_frame = right_box.text_frame
    text_frame.clear()

    paragraph = text_frame.paragraphs[0]
    paragraph.text = "診察時に確認したい点"
    paragraph.font.name = "Yu Gothic"
    paragraph.font.size = Pt(16)
    paragraph.font.bold = True

    questions = [
        "体重低下が意図した範囲か",
        "食欲・疲労・筋力の変化",
        "血糖値・HbA1cとの関連",
        "薬剤変更との関係",
        "心拍・血圧・運動量との関係",
    ]

    for question in questions:
        paragraph = text_frame.add_paragraph()
        paragraph.text = "・" + question
        paragraph.font.name = "Yu Gothic"
        paragraph.font.size = Pt(12)
        paragraph.space_before = Pt(8)

    add_footer(slide)


def add_doctor_memo_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "診察時の確認メモ")

    memo_box = slide.shapes.add_textbox(
        Inches(0.8),
        Inches(1.25),
        Inches(8.4),
        Inches(5.4),
    )
    text_frame = memo_box.text_frame
    text_frame.clear()

    questions = [
        (
            "アブレーション後に体重が継続的に低下しています。"
            "この程度の変化は問題ありませんか。"
        ),
        (
            "食事量や運動量の変化だけで説明できる"
            "範囲でしょうか。"
        ),
        (
            "血糖値、HbA1c、甲状腺機能、栄養状態などの"
            "確認は必要でしょうか。"
        ),
        (
            "現在の薬剤が体重・心拍・血圧に影響している"
            "可能性はありますか。"
        ),
        (
            "今後、どの程度の体重減少や症状があれば、"
            "早めに受診すべきでしょうか。"
        ),
    ]

    for index, question in enumerate(questions, start=1):
        paragraph = (
            text_frame.paragraphs[0]
            if index == 1
            else text_frame.add_paragraph()
        )
        paragraph.text = f"{index}. {question}"
        paragraph.font.name = "Yu Gothic"
        paragraph.font.size = Pt(16)
        paragraph.space_after = Pt(14)

    add_footer(slide)


def make_pptx(
    png_files,
    start_date,
    end_date,
    mode,
    ablation_date,
):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    if mode == "ablation":
        df = load_ablation_data(start_date, end_date)
        summary_rows = calculate_ablation_summary(
            df,
            ablation_date,
        )

        add_overview_slide(
            prs,
            start_date,
            end_date,
            ablation_date,
        )

        add_dashboard_slide(
            prs,
            summary_rows,
            start_date,
            end_date,
            ablation_date,
        )

    for item, png_path in png_files:
        title = (
            f"{ITEMS[item]}（{start_date} ～ {end_date}）"
        )

        subtitle = None
        if mode == "ablation":
            subtitle = (
                f"アブレーション実施日: {ablation_date}。"
                "グラフと前後比較表を併せてご確認ください。"
            )

        add_graph_slide(
            prs,
            title,
            png_path,
            subtitle,
        )

    if mode == "ablation":
        add_doctor_memo_slide(prs)
        pptx_path = ABLATION_PPTX_PATH
    else:
        pptx_path = NORMAL_PPTX_PATH

    prs.save(pptx_path)
    print("PowerPoint作成完了:", pptx_path)

    return pptx_path


def main():
    try:
        parser = argparse.ArgumentParser()

        parser.add_argument(
            "--start-date",
            required=True,
            help="開始日 YYYY-MM-DD",
        )
        parser.add_argument(
            "--end-date",
            required=True,
            help="終了日 YYYY-MM-DD",
        )
        parser.add_argument(
            "--mode",
            choices=["normal", "ablation"],
            default="normal",
            help=(
                "normal: 従来の健康管理グラフ、"
                "ablation: アブレーション前後比較"
            ),
        )
        parser.add_argument(
            "--ablation-date",
            default="2026-05-22",
            help="アブレーション実施日 YYYY-MM-DD",
        )

        args = parser.parse_args()

        start_timestamp = pd.Timestamp(args.start_date)
        end_timestamp = pd.Timestamp(args.end_date)
        ablation_timestamp = pd.Timestamp(
            args.ablation_date
        )

        if start_timestamp > end_timestamp:
            raise ValueError(
                "開始日は終了日以前にしてください。"
            )

        if args.mode == "ablation":
            if not (
                start_timestamp
                <= ablation_timestamp
                <= end_timestamp
            ):
                raise ValueError(
                    "アブレーション日は指定期間内にしてください。"
                )

        png_files = []

        for item in ITEMS.keys():
            run_graph_script(
                args.start_date,
                args.end_date,
                item,
            )
            png_path = export_excel_chart_to_png(item)
            png_files.append((item, png_path))

        pptx_path = make_pptx(
            png_files,
            args.start_date,
            args.end_date,
            args.mode,
            args.ablation_date,
        )

        os.startfile(pptx_path)

    except Exception:
        log = BASE_HEALTH / "make_health_pptx_error.txt"

        with open(log, "w", encoding="utf-8") as file:
            traceback.print_exc(file=file)

        print("エラー内容を保存しました:", log)
        raise


if __name__ == "__main__":
    main()